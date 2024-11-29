import os
import openai
from flask import Flask, render_template, request, redirect, url_for
from werkzeug.utils import secure_filename
from transformers import BertTokenizer, BertForMaskedLM, pipeline
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv

load_dotenv()

openai.api_key = os.getenv("OPENAI_API_KEY")

app = Flask(__name__)

app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'pptx', 'docx'}


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


def extractive_summary(text):
    tokenizer = BertTokenizer.from_pretrained('bert-base-uncased')
    model = BertForMaskedLM.from_pretrained('bert-base-uncased')
    nlp = pipeline("fill-mask", model=model, tokenizer=tokenizer)

    sentences = text.split(".")
    top_sentences = []

    for sentence in sentences:
        mask_sentence = sentence + " [MASK]"
        predictions = nlp(mask_sentence)
        top_sentences.append(predictions[0]['sequence'])

    return " ".join(top_sentences)


def abstractive_summary(text):
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"Summarize the following text in your own words: {text}",
        max_tokens=150
    )
    return response.choices[0].text.strip()


def extract_pdf_text(pdf_path):
    with open(pdf_path, 'rb') as file:
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
    return text


def extract_ppt_text(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text


def extract_word_text(docx_path):
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text


def summarize_document(file_path):
    ext = file_path.split('.')[-1].lower()
    if ext == 'pdf':
        text = extract_pdf_text(file_path)
    elif ext == 'pptx':
        text = extract_ppt_text(file_path)
    elif ext == 'docx':
        text = extract_word_text(file_path)
    else:
        return "Unsupported file format."

    extract_summary = extractive_summary(text)
    abstractive_summary_text = abstractive_summary(extract_summary)

    return abstractive_summary_text


@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        file = request.files['file']
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            summary = summarize_document(filepath)
            return render_template('result.html', summary=summary)

    return render_template('index.html')


if __name__ == "__main__":
    app.run(debug=True)
